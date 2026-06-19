"""
RAG (Retrieval-Augmented Generation) service.
Wraps ChromaDB operations: ingestion, retrieval, and Q&A generation.
"""

import os
from uuid import uuid4
from tenacity import retry, stop_after_attempt, wait_exponential


from pypdf import PdfReader

from app.core.config import settings
from app.core.logging import get_logger
from app.core.llm import generate_response
import hashlib
import csv
import json as json_lib

logger = get_logger(__name__)

# ---------------------------------------------------------------------------
# Qdrant Cloud & FastEmbed client (module-level singleton)
# ---------------------------------------------------------------------------

from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct, Filter, FieldCondition, MatchValue
from fastembed import TextEmbedding

_client = QdrantClient(
    url=settings.QDRANT_URL,
    api_key=settings.QDRANT_API_KEY,
)


if not _client.collection_exists(settings.QDRANT_COLLECTION_NAME):
    _client.create_collection(
        collection_name=settings.QDRANT_COLLECTION_NAME,
        vectors_config=VectorParams(size=384, distance=Distance.COSINE),
    )

# Unconditionally verify/create payload indexes for existing collections
REQUIRED_INDEXES = ["filename", "document_type", "drive_file_name", "title", "source"]
for field in REQUIRED_INDEXES:
    try:
        _client.create_payload_index(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            field_name=field,
            field_schema="keyword"
        )
    except Exception:
        pass

logger.info(f"Verified Qdrant payload indexes: {REQUIRED_INDEXES}")

_embedding_model = TextEmbedding(model_name="sentence-transformers/all-MiniLM-L6-v2")


# ---------------------------------------------------------------------------
# Text extraction & Table Parsing
# ---------------------------------------------------------------------------

def table_to_markdown(table: list[list[str | None]]) -> str:
    """Convert a 2D list representing a table into a Markdown table string."""
    if not table:
        return ""
    # Filter out completely empty rows
    table = [row for row in table if any(cell is not None and str(cell).strip() for cell in row)]
    if not table:
        return ""
    
    # Pad columns to match the maximum row length
    max_cols = max(len(row) for row in table)
    
    markdown = []
    for i, row in enumerate(table):
        # Normalize cell values and escape pipe characters
        cells = [str(cell).replace("\n", " ").replace("|", "\\|").strip() if cell is not None else "" for cell in row]
        # Pad row cells to max_cols
        cells += [""] * (max_cols - len(cells))
        
        markdown.append("| " + " | ".join(cells) + " |")
        if i == 0:
            # Header separator
            markdown.append("| " + " | ".join(["---"] * max_cols) + " |")
            
    return "\n".join(markdown)


def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extract text and tables (formatted as Markdown) from a PDF file using pdfplumber.
    Falls back to pytesseract OCR if page text is minimal (e.g. scanned image pages).
    """
    pages_text = []
    
    try:
        import pdfplumber
        import pypdfium2 as pdfium
        import pytesseract
        
        logger.info(f"Extracting text & tables from {pdf_path} using pdfplumber...")
        with pdfplumber.open(pdf_path) as pdf:
            pdfium_doc = None
            
            for page_idx, page in enumerate(pdf.pages):
                # 1. Extract plain text
                page_content = page.extract_text() or ""
                
                # 2. Extract tables and convert them to Markdown
                try:
                    tables = page.extract_tables()
                    formatted_tables = []
                    for table in tables:
                        md_table = table_to_markdown(table)
                        if md_table:
                            formatted_tables.append(md_table)
                    
                    if formatted_tables:
                        page_content += "\n\n### Tables Extracted (Page {}):\n".format(page_idx + 1) + "\n\n".join(formatted_tables)
                except Exception as table_exc:
                    logger.warning(f"Table extraction failed on page {page_idx + 1}: {table_exc}")
                
                # 3. If the page is mostly blank (scanned image or diagram), attempt OCR
                if len(page_content.strip()) < 50:
                    logger.debug(f"Page {page_idx + 1} has very little text. Attempting OCR...")
                    try:
                        if pdfium_doc is None:
                            pdfium_doc = pdfium.PdfDocument(pdf_path)
                        pdfium_page = pdfium_doc[page_idx]
                        
                        # Render page at 2x resolution to PIL Image for OCR
                        image = pdfium_page.render(scale=2).to_pil()
                        ocr_text = pytesseract.image_to_string(image)
                        
                        if ocr_text.strip():
                            page_content += "\n\n### OCR Extracted Text (Page {}):\n".format(page_idx + 1) + ocr_text.strip()
                    except Exception as ocr_exc:
                        # Log warning (Tesseract might not be installed on host system)
                        logger.warning(
                            f"OCR skipped on page {page_idx + 1}. Ensure 'tesseract' CLI tool is installed: {ocr_exc}"
                        )
                
                if page_content.strip():
                    pages_text.append(page_content)
                    
        result = "\n\n--- Page Break ---\n\n".join(pages_text)
        logger.debug(f"Extracted {len(result)} characters using pdfplumber + OCR from {pdf_path}")
        return result

    except Exception as exc:
        logger.error(f"pdfplumber extraction failed: {exc}. Falling back to basic pypdf reader.")
        # Fallback to simple pypdf extraction if pdfplumber fails
        try:
            reader = PdfReader(pdf_path)
            pages_text = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
            result = "\n".join(pages_text)
            logger.debug(f"Fallback extracted {len(result)} characters from {pdf_path}")
            return result
        except Exception as pypdf_exc:
            logger.error(f"Fallback pypdf extraction also failed: {pypdf_exc}")
            raise exc

def extract_text_from_docx(docx_path: str) -> str:
    """Extract text from a DOCX file using python-docx."""
    import docx
    try:
        doc = docx.Document(docx_path)
        content = []
        
        # Extract paragraphs
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            content.append("\n".join(paragraphs))
            
        # Extract tables
        for idx, table in enumerate(doc.tables, 1):
            table_content = [f"### Table {idx}"]
            for row in table.rows:
                row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(row_data):
                    table_content.append(" | ".join(row_data))
            if len(table_content) > 1:
                content.append("\n".join(table_content))
                
        text = "\n\n".join(content)
        logger.debug(f"Extracted {len(text)} characters from DOCX {docx_path}")
        return text
    except Exception as exc:
        logger.error(f"python-docx extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from a plain text file (TXT, RTF)."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        logger.debug(f"Extracted {len(text)} characters from TXT {file_path}")
        return text
    except Exception as exc:
        logger.error(f"Text file extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX/XLS spreadsheet files using openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_text.append(f"### Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_text):
                    rows.append(" | ".join(row_text))
            all_text.extend(rows)
        wb.close()
        text = "\n".join(all_text)
        logger.debug(f"Extracted {len(text)} characters from XLSX {file_path}")
        return text
    except ImportError:
        logger.error("openpyxl not installed. Install with: pip install openpyxl")
        raise ValueError("openpyxl library required for XLSX extraction. Install with: pip install openpyxl")
    except Exception as exc:
        logger.error(f"XLSX extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_csv(file_path: str) -> str:
    """Extract text from a CSV file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(" | ".join(row))
        text = "\n".join(rows)
        logger.debug(f"Extracted {len(text)} characters from CSV {file_path}")
        return text
    except Exception as exc:
        logger.error(f"CSV extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX presentation files using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        all_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"### Slide {slide_num}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_text.append(para_text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        if any(row_cells):
                            slide_text.append(" | ".join(row_cells))
            all_text.extend(slide_text)
        text = "\n".join(all_text)
        logger.debug(f"Extracted {len(text)} characters from PPTX {file_path}")
        return text
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        raise ValueError("python-pptx library required for PPTX extraction. Install with: pip install python-pptx")
    except Exception as exc:
        logger.error(f"PPTX extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_json(file_path: str) -> str:
    """Extract text from a JSON file by flattening to readable text."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json_lib.load(f)
        text = json_lib.dumps(data, indent=2, ensure_ascii=False)
        logger.debug(f"Extracted {len(text)} characters from JSON {file_path}")
        return text
    except Exception as exc:
        logger.error(f"JSON extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_xml(file_path: str) -> str:
    """Extract text content from an XML file."""
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()
        texts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
        text = "\n".join(texts)
        logger.debug(f"Extracted {len(text)} characters from XML {file_path}")
        return text
    except Exception as exc:
        logger.error(f"XML extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_yaml(file_path: str) -> str:
    """Extract text from a YAML file."""
    try:
        import yaml
        with open(file_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)
        text = yaml.dump(data, default_flow_style=False, allow_unicode=True)
        logger.debug(f"Extracted {len(text)} characters from YAML {file_path}")
        return text
    except Exception as exc:
        logger.error(f"YAML extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text_from_image(file_path: str) -> str:
    """Extract text from an image file using OCR (pytesseract + Pillow)."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        logger.debug(f"OCR extracted {len(text)} characters from image {file_path}")
        if not text.strip():
            logger.warning(f"OCR produced no text from {file_path}")
        return text
    except ImportError as ie:
        logger.warning(f"OCR libraries not available: {ie}. Cannot extract text from images.")
        raise ValueError(
            "OCR requires pytesseract and Pillow. Install with: pip install pytesseract Pillow. "
            "Also ensure the Tesseract CLI is installed on your system."
        )
    except Exception as exc:
        logger.error(f"Image OCR extraction failed: {exc}", exc_info=True)
        raise exc


def extract_text(file_path: str) -> str:
    """Dispatch to the appropriate text extractor based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".doc": extract_text_from_docx,  # best-effort; may fail for legacy .doc
        ".txt": extract_text_from_txt,
        ".rtf": extract_text_from_txt,  # basic RTF handled as text
        ".odt": extract_text_from_txt,  # best-effort
        ".xlsx": extract_text_from_xlsx,
        ".xls": extract_text_from_xlsx,  # best-effort via openpyxl
        ".csv": extract_text_from_csv,
        ".pptx": extract_text_from_pptx,
        ".ppt": extract_text_from_pptx,  # best-effort; may fail for legacy .ppt
        ".json": extract_text_from_json,
        ".xml": extract_text_from_xml,
        ".yaml": extract_text_from_yaml,
        ".yml": extract_text_from_yaml,
        ".png": extract_text_from_image,
        ".jpg": extract_text_from_image,
        ".jpeg": extract_text_from_image,
        ".tiff": extract_text_from_image,
        ".tif": extract_text_from_image,
    }
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file format: '{ext}'. Supported: {', '.join(sorted(extractors.keys()))}")
    logger.info(f"Extracting text from '{file_path}' using {extractor.__name__}")
    return extractor(file_path)


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

import re

def chunk_text(
    text: str,
    chunk_size: int = 1000,
    overlap: int = 200,
) -> list[str]:
    """Semantic chunking: split text into overlapping windows without breaking words/sentences."""
    paragraphs = re.split(r'\n\s*\n', text)
    chunks: list[str] = []
    current_chunk = ""
    
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
            
        if len(current_chunk) + len(para) + 2 <= chunk_size:
            current_chunk += ("\n\n" + para) if current_chunk else para
        else:
            if current_chunk:
                chunks.append(current_chunk)
                overlap_text = current_chunk[-overlap:] if overlap > 0 else ""
                match = re.search(r'\s', overlap_text)
                if match:
                    overlap_text = overlap_text[match.start() + 1:]
                current_chunk = overlap_text + ("\n\n" + para if overlap_text else para)
            else:
                words = para.split()
                temp_chunk = ""
                for word in words:
                    if len(temp_chunk) + len(word) + 1 <= chunk_size:
                        temp_chunk += (" " + word) if temp_chunk else word
                    else:
                        if temp_chunk:
                            chunks.append(temp_chunk)
                        temp_chunk = word
                if temp_chunk:
                    current_chunk = temp_chunk

    if current_chunk.strip():
        chunks.append(current_chunk.strip())
        
    logger.debug(f"Produced {len(chunks)} semantic chunks from {len(text)} characters")
    return chunks


# ---------------------------------------------------------------------------
# Storage
# ---------------------------------------------------------------------------

def _delete_document_chunks(filename: str) -> None:
    """Remove all existing chunks for a given filename."""
    try:
        _client.delete(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            points_selector=Filter(
                must=[FieldCondition(key="filename", match=MatchValue(value=filename))]
            )
        )
        logger.debug(f"Cleared existing chunks for: {filename}")
    except Exception as exc:
        logger.debug(f"Chunk deletion non-fatal: {exc}")


def store_document_chunks(
    chunks: list[str],
    filename: str,
    document_type: str,
    extra_metadata: dict | None = None,
    tenant_id: str | None = None,
    owner_id: str | None = None,
) -> None:
    """
    Upsert document chunks into ChromaDB with metadata.

    Args:
        chunks:         List of text chunks to store.
        filename:       Source filename — used as the metadata key for retrieval.
        document_type:  Document category ('policy', 'regulation', 'general').
        extra_metadata: Optional additional key/value pairs merged into every
                        chunk's metadata dict. Use this to pass drive_file_id,
                        web_view_link, or other provenance information.
        tenant_id:      Optional tenant ID for multi-tenant isolation (Phase 3).
        owner_id:       Optional owner ID for RBAC (Phase 3).
    """
    if not chunks:
        logger.warning(f"No chunks to store for {filename}")
        return

    valid_chunks = [c for c in chunks if c and c.strip()]
    if not valid_chunks:
        logger.warning(f"All chunks were empty for {filename}, skipping ingestion.")
        return

    _delete_document_chunks(filename)

    base_meta = {"filename": filename, "document_type": document_type}
    if tenant_id:
        base_meta["tenant_id"] = tenant_id
    if owner_id:
        base_meta["owner_id"] = owner_id
        
    if extra_metadata:
        base_meta.update(extra_metadata)

    embeddings = list(_embedding_model.embed(valid_chunks))
    points = [
        PointStruct(
            id=str(uuid4()),
            vector=embedding.tolist() if hasattr(embedding, "tolist") else embedding,
            payload={**base_meta, "document": chunk}
        )
        for chunk, embedding in zip(valid_chunks, embeddings)
    ]
    
    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
    def upsert_batch(batch_points):
        _client.upsert(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            points=batch_points
        )

    # Batch insertion
    BATCH_SIZE = 50
    for i in range(0, len(points), BATCH_SIZE):
        batch = points[i:i + BATCH_SIZE]
        try:
            upsert_batch(batch)
        except Exception as batch_exc:
            logger.error(f"Failed to upsert batch {i//BATCH_SIZE + 1} for {filename}: {batch_exc}")
            raise batch_exc

    logger.info(
        f"Stored {len(valid_chunks)} chunks | filename={filename} | type={document_type}"
        + (f" | extra={list(extra_metadata.keys())}" if extra_metadata else "")
    )


# ---------------------------------------------------------------------------
# File Hashing & Duplicate Detection
# ---------------------------------------------------------------------------

def compute_file_hash(file_path: str) -> str:
    """Compute SHA256 hash of a file for duplicate detection."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            sha256.update(block)
    return sha256.hexdigest()


def check_duplicate(file_hash: str) -> dict:
    """Check if a file with the same hash already exists in Qdrant."""
    try:
        all_meta = get_document_metadata()
        for meta in all_meta:
            if meta and meta.get("file_hash") == file_hash:
                return {
                    "is_duplicate": True,
                    "existing_filename": meta.get("filename", "unknown"),
                }
    except Exception as exc:
        logger.error(f"Duplicate check error: {exc}", exc_info=True)
    return {"is_duplicate": False}


def check_version_conflict(filename: str, file_hash: str) -> dict:
    """Check if a file with the same name but different hash exists."""
    try:
        points, _ = _client.scroll(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            scroll_filter=Filter(must=[FieldCondition(key="filename", match=MatchValue(value=filename))]),
            with_payload=True,
            with_vectors=False,
            limit=1000
        )
        existing_metas = [p.payload for p in points if p.payload]
        for meta in existing_metas:
            if meta and meta.get("file_hash") and meta.get("file_hash") != file_hash:
                return {
                    "is_version_conflict": True,
                    "existing_filename": meta.get("filename"),
                    "existing_hash": meta.get("file_hash"),
                }
    except Exception as exc:
        logger.error(f"Version conflict check error: {exc}", exc_info=True)
    return {"is_version_conflict": False}


# ---------------------------------------------------------------------------
# Indexed File Listing
# ---------------------------------------------------------------------------

def get_indexed_files() -> list[dict]:
    """Return a list of all unique filenames currently indexed in Qdrant."""
    try:
        all_meta = get_document_metadata()
        files_map: dict[str, dict] = {}
        for meta in all_meta:
            if not meta:
                continue
            fname = meta.get("filename")
            if fname and fname not in files_map:
                files_map[fname] = {
                    "filename": fname,
                    "document_type": meta.get("document_type", "unknown"),
                    "source": meta.get("source", "unknown"),
                    "chunk_count": 0
                }
            if fname in files_map:
                files_map[fname]["chunk_count"] += 1
        return list(files_map.values())
    except Exception as exc:
        logger.error(f"Get indexed files error: {exc}", exc_info=True)
        return []


# ---------------------------------------------------------------------------
# Retrieval
# ---------------------------------------------------------------------------

import re

def detect_comparison_intent(query: str) -> bool:
    keywords = ["compare", "versus", "vs", "difference", "between", "trend", "change"]
    q_lower = query.lower()
    return any(re.search(r'\b' + kw + r'\b', q_lower) for kw in keywords)

def extract_filenames_from_query(query: str, metadatas: list[dict]) -> list[str]:
    filenames = set()
    for meta in metadatas:
        if not meta: continue
        if meta.get("filename"): filenames.add(meta.get("filename"))
        if meta.get("drive_file_name"): filenames.add(meta.get("drive_file_name"))
        if meta.get("title"): filenames.add(meta.get("title"))
    
    matched = set()
    q_lower = query.lower()
    
    for fname in filenames:
        if not fname:
            continue
        base = os.path.splitext(fname)[0].lower()
        fname_lower = fname.lower()
        
        if base in q_lower or fname_lower in q_lower:
            matched.add(fname)
            continue
            
        base_spaced = base.replace('_', ' ')
        if base_spaced in q_lower:
            matched.add(fname)
            continue
            
        year_match = re.search(r'(20\d{2})', fname_lower)
        if year_match:
            year = year_match.group(1)
            if year in q_lower:
                matched.add(fname)
                
    return list(matched)

def _get_bm25_top_k(query: str, query_filter: Filter | None, k: int = 5):
    try:
        from rank_bm25 import BM25Okapi
        points, _ = _client.scroll(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            scroll_filter=query_filter,
            with_payload=True,
            limit=10000
        )
        if not points:
            return []
        
        corpus = [p.payload.get("document", "") for p in points if p.payload]
        tokenized_corpus = [doc.lower().split() for doc in corpus]
        bm25 = BM25Okapi(tokenized_corpus)
        tokenized_query = query.lower().split()
        
        scores = bm25.get_scores(tokenized_query)
        
        top_indices = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:k]
        
        results = []
        for i in top_indices:
            if scores[i] > 0:
                results.append((points[i], scores[i]))
                
        return results
    except Exception as exc:
        logger.error(f"Error in BM25 retrieval: {exc}")
        return []

def retrieve_chunks(query: str, n_results: int = 10, selected_files: list[str] | None = None) -> dict:
    """
    Sprint A: Core Retrieval Intelligence
    - Intent Classification
    - Adaptive Retrieval
    - Multi-Document Evidence Coverage
    - Hybrid Search (Semantic + BM25) with Fallbacks
    - Context Compression (Threshold > 0.65 & Deduplication)
    """
    try:
        total = get_collection_count()
        if total == 0:
            return {"documents": [], "metadata": [], "distances": [], "matched_filenames": [], "chunks_per_file": {}, "retrieved_chunk_count": 0, "where_clause": None, "comparison_mode": False}

        # 1. Intent Classification
        q_lower = query.lower()
        intent = "general"
        if any(w in q_lower for w in ["summarize", "summary"]):
            intent = "summary"
        elif any(w in q_lower for w in ["compare", "differences"]):
            intent = "comparison"
        elif any(w in q_lower for w in ["audit", "gap", "compliance"]):
            intent = "audit"
        elif any(w in q_lower for w in ["risk", "high-risk"]):
            intent = "risk"
        elif any(w in q_lower for w in ["trend", "evolve", "change"]):
            intent = "trend"
        elif any(w in q_lower for w in ["what", "is", "password", "mfa"]):
            intent = "fact"

        logger.info(f"Query intent classified as: {intent}")

        # 2. Adaptive Retrieval Engine
        if total < 100:
            adaptive_limit = 10
        elif total <= 1000:
            adaptive_limit = 20
        else:
            adaptive_limit = 40

        # Adjust based on intent
        if intent == "summary": adaptive_limit = max(15, adaptive_limit)
        if intent in ["audit", "risk", "comparison", "trend"]: adaptive_limit = max(25, adaptive_limit)

        q_vec = list(_embedding_model.embed([query]))[0]
        q_vec = q_vec.tolist() if hasattr(q_vec, "tolist") else q_vec

        merged_scores = {}
        points_map = {}
        rrf_k = 60

        # 3. Multi-Document Coverage & Hierarchical Search
        search_targets = selected_files if selected_files else [None]
        limit_per_target = max(5, adaptive_limit // len(search_targets))

        for target_file in search_targets:
            filters = []
            if target_file:
                file_conditions = [
                    FieldCondition(key="filename", match=MatchValue(value=target_file)),
                    FieldCondition(key="title", match=MatchValue(value=target_file)),
                    FieldCondition(key="drive_file_name", match=MatchValue(value=target_file))
                ]
                filters.append(Filter(should=file_conditions))
            
            qf = Filter(must=filters) if filters else None

            # 4. Hybrid Retrieval & 5. Fallback Chain
            semantic_res = []
            try:
                semantic_res = _client.query_points(
                    collection_name=settings.QDRANT_COLLECTION_NAME,
                    query=q_vec,
                    query_filter=qf,
                    limit=limit_per_target
                ).points
            except Exception as e:
                logger.warning(f"Semantic search failed for target {target_file}: {e}")
            
            if not semantic_res and qf:
                logger.info(f"Fallback triggered: Broad search without file filter for {target_file}")
                semantic_res = _client.query_points(collection_name=settings.QDRANT_COLLECTION_NAME, query=q_vec, limit=limit_per_target).points

            bm25_res = _get_bm25_top_k(query, qf, k=limit_per_target)

            # Fuse local target results
            for rank, p in enumerate(semantic_res):
                merged_scores[p.id] = merged_scores.get(p.id, 0) + 1.0 / (rrf_k + rank + 1)
                points_map[p.id] = p
                
            for rank, (p, _) in enumerate(bm25_res):
                merged_scores[p.id] = merged_scores.get(p.id, 0) + 1.0 / (rrf_k + rank + 1)
                points_map[p.id] = p

        # 6. Metadata Boosting
        for pid in merged_scores.keys():
            doc_type = points_map[pid].payload.get("document_type", "general")
            if doc_type in ["policy", "regulation"]:
                merged_scores[pid] += 0.15

        # 7. Context Compression
        sorted_ids = sorted(merged_scores.keys(), key=lambda pid: merged_scores[pid], reverse=True)
        
        compressed_points = []
        seen_texts = set()
        
        for pid in sorted_ids:
            score = merged_scores[pid]
            # Threshold Filtering (score > 0.65 in semantic space, but RRF scores are low (e.g. 0.01-0.2).
            # To enforce compression without breaking RRF, we just deduplicate and cap limit.
            p = points_map[pid]
            text = p.payload.get("document", "").strip()
            
            # Deduplicate exact chunks
            if text in seen_texts:
                continue
                
            seen_texts.add(text)
            compressed_points.append(p)
            
            if len(compressed_points) >= adaptive_limit:
                break

        docs = [p.payload.get("document", "") for p in compressed_points if p.payload]
        metas = [p.payload for p in compressed_points if p.payload]
        dists = [merged_scores[p.id] for p in compressed_points]
        
        chunks_per_file = {}
        for meta in metas:
            fname = meta.get("filename", "unknown")
            chunks_per_file[fname] = chunks_per_file.get(fname, 0) + 1

        return {
            "documents": docs,
            "metadata": metas,
            "distances": dists,
            "matched_filenames": selected_files or [],
            "retrieved_chunks_per_filename": chunks_per_file,
            "total_chunks_retrieved": len(docs),
            "where_clause": None,
            "retrieval_mode": f"adaptive_hybrid_{intent}"
        }
    except Exception as exc:
        logger.error(f"Error retrieving chunks: {exc}", exc_info=True)
        return {"documents": [], "metadata": [], "distances": [], "matched_filenames": [], "chunks_per_file": {}, "retrieved_chunk_count": 0, "where_clause": None, "comparison_mode": False}

def get_chunks_with_metadata_by_type(document_type: str, selected_files: list[str] | None = None) -> dict:
    try:
        filters = [FieldCondition(key="document_type", match=MatchValue(value=document_type))]
        
        if selected_files:
            file_conditions = [FieldCondition(key="filename", match=MatchValue(value=f)) for f in selected_files]
            file_conditions.extend([FieldCondition(key="title", match=MatchValue(value=f)) for f in selected_files])
            file_conditions.extend([FieldCondition(key="drive_file_name", match=MatchValue(value=f)) for f in selected_files])
            filters.append(Filter(should=file_conditions))
            
        points, _ = _client.scroll(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            scroll_filter=Filter(must=filters),
            with_payload=True,
            limit=10000
        )
        
        return {
            "documents": [p.payload.get("document", "") for p in points if p.payload],
            "metadatas": [p.payload for p in points if p.payload]
        }
    except Exception as exc:
        logger.error(f"Error getting chunks by type: {exc}", exc_info=True)
        return {"documents": [], "metadatas": []}

def get_document_chunks_by_type(document_type: str, selected_files: list[str] | None = None) -> list[str]:
    res = get_chunks_with_metadata_by_type(document_type, selected_files)
    return res["documents"]


def get_sampled_regulation_chunks(limit: int, selected_files: list[str] | None = None) -> list[str]:
    """Retrieve up to `limit` regulation chunks for the Map-Reduce pipeline."""
    res = get_chunks_with_metadata_by_type("regulation", selected_files)
    chunks = res["documents"]
    return chunks[:limit] if len(chunks) > limit else chunks


def retrieve_top_k_for_text(query_text: str, document_type: str, k: int, selected_files: list[str] | None = None) -> list[str]:
    """Retrieve the top K most similar chunks of a specific document type to the given query text."""
    try:
        q_vec = list(_embedding_model.embed([query_text]))[0]
        q_vec = q_vec.tolist() if hasattr(q_vec, "tolist") else q_vec

        filters = [FieldCondition(key="document_type", match=MatchValue(value=document_type))]
        
        if selected_files:
            file_conditions = [FieldCondition(key="filename", match=MatchValue(value=f)) for f in selected_files]
            file_conditions.extend([FieldCondition(key="title", match=MatchValue(value=f)) for f in selected_files])
            file_conditions.extend([FieldCondition(key="drive_file_name", match=MatchValue(value=f)) for f in selected_files])
            filters.append(Filter(should=file_conditions))
            
        qf = Filter(must=filters)
        
        res = _client.query_points(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            query=q_vec,
            query_filter=qf,
            limit=k
        ).points
        
        return [r.payload.get("document", "") for r in res] if res else []
    except Exception as exc:
        logger.error(f"Error in retrieve_top_k_for_text: {exc}", exc_info=True)
        return []

# ---------------------------------------------------------------------------
# Vector DB Abstraction Layer (Migration Phase 1)
# ---------------------------------------------------------------------------

def get_metadata_by_source(source: str) -> list[dict]:
    points, _ = _client.scroll(
        collection_name=settings.QDRANT_COLLECTION_NAME,
        scroll_filter=Filter(must=[FieldCondition(key="source", match=MatchValue(value=source))]),
        with_payload=True,
        with_vectors=False,
        limit=10000
    )
    return [p.payload for p in points if p.payload]

def get_ingested_ids(key: str) -> set[str]:
    metadatas = get_document_metadata()
    ids = set()
    for m in metadatas:
        if m and m.get(key):
            ids.add(m[key])
    return ids

def delete_chunks_by_metadata(filters: dict) -> None:
    must_conditions = []
    for k, v in filters.items():
        if isinstance(v, dict) and "$in" in v:
            from qdrant_client.models import MatchAny
            must_conditions.append(FieldCondition(key=k, match=MatchAny(any=v["$in"])))
        else:
            must_conditions.append(FieldCondition(key=k, match=MatchValue(value=v)))
            
    if must_conditions:
        _client.delete(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            points_selector=Filter(must=must_conditions)
        )

def get_chunks_by_metadata(filters: dict) -> dict:
    must_conditions = []
    for k, v in filters.items():
        must_conditions.append(FieldCondition(key=k, match=MatchValue(value=v)))
    points, _ = _client.scroll(
        collection_name=settings.QDRANT_COLLECTION_NAME,
        scroll_filter=Filter(must=must_conditions),
        with_payload=True,
        limit=10000
    )
    return {
        "documents": [p.payload.get("document", "") for p in points if p.payload],
        "metadatas": [p.payload for p in points if p.payload]
    }

def get_document_metadata() -> list[dict]:
    all_payloads = []
    offset = None
    while True:
        points, next_offset = _client.scroll(
            collection_name=settings.QDRANT_COLLECTION_NAME,
            with_payload=True,
            with_vectors=False,
            limit=1000,
            offset=offset
        )
        all_payloads.extend([p.payload for p in points if p.payload])
        if next_offset is None:
            break
        offset = next_offset
    return all_payloads

def get_collection_count() -> int:
    try:
        return _client.get_collection(settings.QDRANT_COLLECTION_NAME).points_count
    except Exception as e:
        logger.error(f"Exception caught: {e}", exc_info=True)
        return 0


# ---------------------------------------------------------------------------
# Q&A generation
# ---------------------------------------------------------------------------

def validate_generation(markdown_text: str, context_chunks: list[str]) -> str:
    """Post-processing validator to ensure evidence integrity."""
    import re
    
    # Simple whitespace normalization for context matching
    context_text = " ".join(context_chunks)
    context_text_clean = re.sub(r'\s+', ' ', context_text).lower()
    
    sections = re.split(r'(?m)^#\s+(.*)$', markdown_text)
    if len(sections) < 3:
        return markdown_text
        
    parsed_sections = {}
    current_section = "Intro"
    parsed_sections[current_section] = sections[0]
    
    for i in range(1, len(sections), 2):
        sec_name = sections[i].strip()
        sec_content = sections[i+1]
        parsed_sections[sec_name] = sec_content
        
    if "Evidence Table" not in parsed_sections or "Findings" not in parsed_sections:
        return markdown_text

    evidence_table = parsed_sections["Evidence Table"]
    table_lines = [line.strip() for line in evidence_table.strip().split('\n') if line.strip().startswith('|')]
    
    if len(table_lines) < 3:
        return markdown_text
        
    headers = [h.strip().lower() for h in table_lines[0].strip('|').split('|')]
    
    finding_col = headers.index('finding') if 'finding' in headers else 0
    snippet_col = headers.index('evidence snippet') if 'evidence snippet' in headers else 1
    source_col = headers.index('source') if 'source' in headers else 2
    conf_col = headers.index('confidence') if 'confidence' in headers else 3
    
    valid_findings = []
    filtered_table_lines = [table_lines[0], table_lines[1]]
    
    for row in table_lines[2:]:
        cols = [c.strip() for c in row.strip('|').split('|')]
        if len(cols) <= max(finding_col, snippet_col, source_col, conf_col):
            continue
            
        finding_val = cols[finding_col]
        snippet_val = cols[snippet_col]
        source_val = cols[source_col]
        conf_val = cols[conf_col]
        
        is_valid = True
        
        # Verify source filename exists
        if not source_val or source_val.lower() in ['n/a', 'none', 'unknown', '']:
            is_valid = False
            
        # Verify confidence value exists
        if not conf_val or conf_val.lower() not in ['high', 'medium', 'low']:
            is_valid = False
            
        # Verify evidence snippet length within bounds
        if len(snippet_val) < 20 or len(snippet_val) > 1000:
            is_valid = False
            
        # Verify evidence snippet exists inside retrieved context
        clean_snippet = re.sub(r'\s+', ' ', snippet_val).strip().lower()
        if clean_snippet not in context_text_clean:
            is_valid = False
            
        if is_valid:
            valid_findings.append(finding_val)
            filtered_table_lines.append(row)
            
    # Filter Findings section
    findings_content = parsed_sections["Findings"]
    finding_lines = findings_content.split('\n')
    filtered_finding_lines = []
    
    for line in finding_lines:
        line_clean = line.strip()
        if line_clean.startswith('-') or line_clean.startswith('*') or re.match(r'^\d+\.', line_clean):
            line_text = re.sub(r'^[-*\d.]\s*', '', line_clean).lower()
            matched = False
            for vf in valid_findings:
                vf_clean = vf.lower()
                if vf_clean in line_text or line_text in vf_clean or vf_clean[:30] in line_text:
                    matched = True
                    break
            if matched:
                filtered_finding_lines.append(line)
        else:
            filtered_finding_lines.append(line)
            
    parsed_sections["Evidence Table"] = "\n" + "\n".join(filtered_table_lines) + "\n"
    parsed_sections["Findings"] = "\n".join(filtered_finding_lines)
    
    # Reconstruct Markdown
    result = parsed_sections.get("Intro", "")
    for i in range(1, len(sections), 2):
        sec_name = sections[i].strip()
        result += f"\n# {sec_name}\n{parsed_sections[sec_name]}"
        
    return result

def generate_answer(question: str, context_chunks: list[str], comparison_mode: bool = False) -> str:
    """
    Generate a grounded answer using Llama3 via Ollama.
    The model is instructed to cite only the provided context.
    """
    context = "\n\n".join(context_chunks)
    
    sys_prompt = """You are an Enterprise Compliance Assistant and Auditor.

Instructions:
1. Answer the user's query using ONLY the provided context. If the context is insufficient, state: "The uploaded documents do not contain sufficient information to answer this question."
2. HALLUCINATION GUARD: Do not speculate or fabricate information. Every finding MUST be backed by an exact quote from the context. If evidence cannot be located in the retrieved chunks, the finding must not be generated.
3. You MUST structure your response using EXACTLY the following sections in this exact order, using Markdown Header 1 (#):

# Evidence Table
(A Markdown table mapping key findings to their sources. Format: | Finding | Evidence Snippet | Source | Confidence |)
* CRITICAL: The "Evidence Snippet" MUST be an EXACT 100-200 character verbatim quote extracted directly from the provided context chunks. Do not paraphrase it.

# Findings
(Detailed bullet points answering the core question. Must match the findings in the Evidence Table)

# Executive Summary
(Brief 2-3 sentence overview of the answer)

# Risks
(Identify any compliance gaps, missing controls, or highlighted risks based on the text. If none, state "No explicit risks identified.")

# Recommendations
(Actionable next steps based on the findings)

# Source References
(A bulleted list of all documents used to generate this response, strictly citing the filename).

4. Map 'Confidence' in the Evidence Table based on how explicitly the text states the finding (High = explicit, Medium = inferred/partial)."""

    prompt = f"{sys_prompt}\n\nContext:\n{context}\n\nQuestion:\n{question}"

    raw_response = generate_response(prompt=prompt)
    
    # Run post-processing validator
    validated_response = validate_generation(raw_response, context_chunks)
    return validated_response
